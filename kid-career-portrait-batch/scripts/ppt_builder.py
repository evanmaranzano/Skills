#!/usr/bin/env python3
"""Create comparison PPTX for kid-career-portrait-batch."""

from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
import tempfile
import shutil
import xml.etree.ElementTree as ET

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
except ImportError:
    Image = None
    Presentation = None
    RGBColor = None
    MSO_SHAPE = None
    Inches = None
    Pt = None


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
ET.register_namespace("p", P_NS)
ET.register_namespace("a", "http://schemas.openxmlformats.org/drawingml/2006/main")
ET.register_namespace("r", "http://schemas.openxmlformats.org/officeDocument/2006/relationships")


def _make_cover_image(image_path: Path, target_ratio: float, asset_dir: Path, stem: str) -> Path:
    """Create a center-cropped image that fills the target ratio without letterboxing."""
    asset_dir.mkdir(parents=True, exist_ok=True)
    out_path = asset_dir / f"{stem}.png"
    with Image.open(image_path) as img:
        img = img.convert("RGB")
        w, h = img.size
        src_ratio = w / h
        if src_ratio > target_ratio:
            new_w = round(h * target_ratio)
            left = max(0, (w - new_w) // 2)
            img = img.crop((left, 0, left + new_w, h))
        elif src_ratio < target_ratio:
            new_h = round(w / target_ratio)
            top = max(0, (h - new_h) // 2)
            img = img.crop((0, top, w, top + new_h))
        img.save(out_path, format="PNG")
    return out_path


def _add_picture_cover(slide, image_path: Path, left, top, width, height, asset_dir: Path, stem: str):
    cover = _make_cover_image(image_path, float(width) / float(height), asset_dir, stem)
    return slide.shapes.add_picture(str(cover), left, top, width=width, height=height)


def _fit_image(image_path: Path, box_left, box_top, box_width, box_height):
    with Image.open(image_path) as img:
        w, h = img.size
    scale = min(box_width / w, box_height / h)
    width = int(w * scale)
    height = int(h * scale)
    left = int(box_left + (box_width - width) / 2)
    top = int(box_top + (box_height - height) / 2)
    return left, top, width, height


def _add_picture_contain(slide, image_path: Path, left, top, width, height):
    pic_left, pic_top, pic_w, pic_h = _fit_image(image_path, left, top, width, height)
    return slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_w, height=pic_h)


def _add_label(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.15), Inches(6.5), Inches(0.35))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(80, 80, 80)
    return box


def _add_caption(slide, text: str, left, top, width):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(90, 90, 90)
    return box


def _add_border(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    shape.line.color.rgb = RGBColor(225, 225, 225)
    shape.line.width = Pt(1)
    return shape


def _timing_xml(spid: int) -> ET.Element:
    xml = f"""
<p:timing xmlns:p=\"{P_NS}\">
  <p:tnLst>
    <p:par>
      <p:cTn id=\"1\" dur=\"indefinite\" restart=\"never\" nodeType=\"tmRoot\">
        <p:childTnLst>
          <p:seq concurrent=\"1\" nextAc=\"seek\">
            <p:cTn id=\"2\" dur=\"indefinite\" nodeType=\"mainSeq\">
              <p:childTnLst>
                <p:par>
                  <p:cTn id=\"3\" fill=\"hold\">
                    <p:stCondLst><p:cond delay=\"indefinite\"/></p:stCondLst>
                    <p:childTnLst>
                      <p:set>
                        <p:cBhvr>
                          <p:cTn id=\"4\" dur=\"1\" fill=\"hold\"/>
                          <p:tgtEl><p:spTgt spid=\"{spid}\"/></p:tgtEl>
                          <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                        </p:cBhvr>
                        <p:to><p:strVal val=\"visible\"/></p:to>
                      </p:set>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst><p:cond evt=\"onPrev\" delay=\"0\"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
            <p:nextCondLst><p:cond evt=\"onNext\" delay=\"0\"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst><p:bldP spid=\"{spid}\" grpId=\"0\"/></p:bldLst>
</p:timing>
"""
    return ET.fromstring(xml)


def _patch_slide_animation(pptx_path: Path, slide_index_to_spid: dict[int, int]) -> None:
    if not slide_index_to_spid:
        return
    with tempfile.TemporaryDirectory() as td:
        td_path = Path(td)
        with ZipFile(pptx_path, "r") as zin:
            zin.extractall(td_path)
        for slide_idx, spid in slide_index_to_spid.items():
            slide_xml = td_path / "ppt" / "slides" / f"slide{slide_idx}.xml"
            if not slide_xml.exists():
                continue
            root = ET.parse(slide_xml).getroot()
            for old in list(root.findall(f"{{{P_NS}}}timing")):
                root.remove(old)
            for c_nv_pr in root.findall(f".//{{{P_NS}}}cNvPr"):
                if c_nv_pr.get("id") == str(spid):
                    c_nv_pr.set("hidden", "1")
            timing = _timing_xml(spid)
            ext = root.find(f"{{{P_NS}}}extLst")
            if ext is not None:
                root.insert(list(root).index(ext), timing)
            else:
                root.append(timing)
            ET.ElementTree(root).write(slide_xml, encoding="utf-8", xml_declaration=True)
        tmp_out = pptx_path.with_suffix(".tmp.pptx")
        with ZipFile(tmp_out, "w", ZIP_DEFLATED) as zout:
            for file in td_path.rglob("*"):
                if file.is_file():
                    zout.write(file, file.relative_to(td_path).as_posix())
        shutil.move(str(tmp_out), str(pptx_path))


def build_comparison_ppt(records: list[dict], output_dir: Path, ppt_name: str = "kid-career-portraits.pptx") -> Path | None:
    """Build a PPTX where right-side future portrait appears on click."""
    if Presentation is None or Image is None:
        raise RuntimeError("PPTX generation requires dependencies: pip install -r requirements.txt")
    valid = []
    for r in records:
        current = Path(r.get("input_file", ""))
        future = Path(r.get("output_file", ""))
        if r.get("status") in {"success", "skipped"} and current.exists() and future.exists():
            valid.append((r, current, future))
    if not valid:
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    animation_targets = {}

    # Portrait frames. Ratio is 3:4, matching the default 1536x2048 image size.
    left_box = (Inches(0.95), Inches(0.78), Inches(4.65), Inches(6.2))
    right_box = (Inches(7.73), Inches(0.78), Inches(4.65), Inches(6.2))
    asset_dir = output_dir / "ppt_assets"

    for idx, (record, current, future) in enumerate(valid, 1):
        slide = prs.slides.add_slide(blank)
        name = record.get("name", "")
        career = record.get("career", "")
        _add_label(slide, f"{name}｜理想职业：{career}")
        _add_border(slide, *left_box)
        _add_border(slide, *right_box)
        _add_caption(slide, "现在", left_box[0], Inches(6.95), left_box[2])
        _add_caption(slide, "未来职业照（按空格显示）", right_box[0], Inches(6.95), right_box[2])
        _add_picture_cover(slide, current, *left_box, asset_dir, f"slide_{idx:03d}_current")
        future_pic = _add_picture_contain(slide, future, *right_box)
        animation_targets[idx] = future_pic.shape_id

    pptx_path = output_dir / ppt_name
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)
    _patch_slide_animation(pptx_path, animation_targets)
    return pptx_path
